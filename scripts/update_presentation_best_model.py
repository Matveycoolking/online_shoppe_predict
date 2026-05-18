from pathlib import Path
from shutil import copy2

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE_PPTX = ROOT / "presentation_updated_preview_source.pptx"
PROJECT_PPTX = ROOT / "Презентация_Покупательское_намерение_Матвей.pptx"
DOWNLOADS_PPTX = (
    Path.home()
    / "Downloads"
    / "Презентация_Покупательское_намерение_Матвей.pptx"
)


def set_shape_text(slide, shape_index: int, text: str) -> None:
    shape = slide.shapes[shape_index]
    if not hasattr(shape, "text_frame"):
        raise ValueError(f"Shape {shape_index} has no text frame")
    shape.text_frame.clear()
    shape.text_frame.paragraphs[0].text = text


def main() -> None:
    if not SOURCE_PPTX.exists():
        raise FileNotFoundError(SOURCE_PPTX)

    copy2(SOURCE_PPTX, PROJECT_PPTX)
    prs = Presentation(PROJECT_PPTX)

    # Slide 5: methods.
    slide = prs.slides[4]
    set_shape_text(slide, 8, "Итоговая модель")
    set_shape_text(
        slide,
        9,
        'Random Forest Tuned\nRandomizedSearchCV\nclass_weight="balanced"\nЛучший F1 = 0.6615',
    )
    set_shape_text(
        slide,
        10,
        "Все модели обучались внутри sklearn Pipeline: feature engineering + preprocessing + estimator. "
        "Это позволяет backend принимать сырые пользовательские данные.",
    )

    # Slide 6: comparison table. Keep the slide compact: top-3 useful rows.
    slide = prs.slides[5]
    values = {
        11: "Random Forest Tuned",
        12: "0.893",
        13: "0.648",
        14: "0.675",
        15: "0.662",
        16: "0.923",
        17: "CatBoost",
        18: "0.902",
        19: "0.727",
        20: "0.586",
        21: "0.649",
        22: "0.929",
        23: "Logistic Regression",
        24: "0.859",
        25: "0.529",
        26: "0.809",
        27: "0.640",
        28: "0.920",
        29: (
            "Random Forest Tuned выбран по F1-score.\n"
            "CatBoost лучше по ROC-AUC, но уступил по основной метрике.\n"
            "Logistic Regression полезна как baseline и даёт максимальный recall."
        ),
    }
    for shape_index, text in values.items():
        set_shape_text(slide, shape_index, text)

    # Slide 7: final model quality and confusion matrix.
    slide = prs.slides[6]
    values = {
        5: "Random Forest Tuned",
        6: "0.893",
        8: "0.662",
        10: "0.923",
        12: "0.675",
        14: "Истинно отрицательный\n1944",
        15: "Ложно положительный\n140",
        16: "Ложно отрицательный\n124",
        17: "Истинно положительный\n258",
    }
    for shape_index, text in values.items():
        set_shape_text(slide, shape_index, text)

    # Slide 10: final result.
    slide = prs.slides[9]
    values = {
        5: (
            "Реализован полный ML-процесс от EDA до веб-демонстрации.\n"
            "Обучены и сравнены четыре модели.\n"
            "Лучшая модель — Random Forest Tuned по F1-score.\n"
            "Backend принимает сырые данные и возвращает вероятность покупки.\n"
            "Реализован Frontend."
        ),
        6: "0.6615",
        8: "0.9232",
        10: "4",
    }
    for shape_index, text in values.items():
        set_shape_text(slide, shape_index, text)

    prs.save(PROJECT_PPTX)
    if DOWNLOADS_PPTX.parent.exists():
        copy2(PROJECT_PPTX, DOWNLOADS_PPTX)

    print(PROJECT_PPTX)
    print(DOWNLOADS_PPTX)


if __name__ == "__main__":
    main()
