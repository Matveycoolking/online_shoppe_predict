from __future__ import annotations

from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "Презентация_Покупательское_намерение_Матвей.pptx"
DOWNLOADS_COPY = (
    Path.home()
    / "Downloads"
    / "Презентация_Покупательское_намерение_Матвей.pptx"
)

DARK = RGBColor(67, 62, 76)
GREEN = RGBColor(137, 196, 0)
CYAN = RGBColor(50, 194, 215)
LIGHT = RGBColor(245, 248, 250)
TEXT = RGBColor(46, 46, 52)
WHITE = RGBColor(255, 255, 255)


def clear_text(shape) -> None:
    shape.text_frame.clear()


def set_text(shape, lines: list[str], font_size: int = 17, bold_first: bool = False) -> None:
    clear_text(shape)
    frame = shape.text_frame
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(4)
        if index == 0 and bold_first:
            paragraph.font.bold = True
            paragraph.font.color.rgb = DARK


def set_center_text(shape, text: str, font_size: int = 18, bold: bool = True) -> None:
    clear_text(shape)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = WHITE if bold else TEXT


def add_card(slide, left, top, width, height, title: str, body: list[str], accent=GREEN) -> None:
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), Inches(0.35))
    set_text(title_box, [title], font_size=17, bold_first=True)

    body_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.55), width - Inches(0.32), height - Inches(0.65))
    set_text(body_box, body, font_size=14)


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def update_slide_7(prs: Presentation) -> None:
    slide = prs.slides[6]
    set_text(
        slide.shapes[4],
        [
            "Схема проверки",
            "Train/test: 80/20 со stratify по Revenue.",
            "Preprocessing обучается только на train внутри Pipeline.",
            "Основной критерий выбора: F1-score.",
        ],
        font_size=16,
        bold_first=True,
    )

    add_card(
        slide,
        Inches(0.93),
        Inches(3.42),
        Inches(4.45),
        Inches(0.92),
        "Cross-validation",
        [
            "F1 = 0.685 ± 0.016",
            "ROC-AUC = 0.930 ± 0.004",
            "Качество стабильно на разных разбиениях.",
        ],
        accent=CYAN,
    )

    add_card(
        slide,
        Inches(5.52),
        Inches(5.18),
        Inches(4.93),
        Inches(0.95),
        "Переобучение",
        [
            "Train F1 = 0.947, test F1 = 0.673.",
            "Есть признаки переобучения, но test/CV остаются рабочими.",
        ],
        accent=GREEN,
    )


def update_slide_8(prs: Presentation) -> None:
    slide = prs.slides[7]
    set_text(
        slide.shapes[6],
        [
            "Top feature importance",
            "PageValues — 0.239",
            "HasPageValue — 0.187",
            "ExitRates — 0.059",
            "TotalDuration — 0.051",
            "ProductRelated_Duration — 0.045",
        ],
        font_size=16,
        bold_first=True,
    )

    set_center_text(slide.shapes[7], "Анализ ошибок", font_size=18)
    set_text(
        slide.shapes[8],
        [
            "False Negative: реальные покупатели, которых модель не нашла.",
            "False Positive: лишние маркетинговые действия.",
            "Для бизнеса важнее контролировать пропущенные покупки.",
        ],
        font_size=15,
    )

    set_center_text(slide.shapes[9], "Практический вывод", font_size=18)
    set_text(
        slide.shapes[10],
        [
            "Модель опирается на поведение пользователя: ценность страниц, выходы, длительность и активность на товарах.",
            "Прогноз лучше использовать как дополнительный сигнал для персонализации.",
        ],
        font_size=15,
    )

    set_center_text(
        slide.shapes[11],
        "Интерпретация показывает: важнее всего не технические поля, а действия пользователя в сессии.",
        font_size=12,
        bold=False,
    )


def update_slide_9(prs: Presentation) -> None:
    slide = prs.slides[8]
    set_center_text(slide.shapes[7], "Docker + Frontend", font_size=17)

    set_text(
        slide.shapes[11],
        [
            "Инференс",
            "Backend endpoint: POST /predict.",
            "Модель сохранена как preprocessing + estimator pipeline.",
            "API возвращает prediction и purchase_probability.",
            "Swagger: http://127.0.0.1:8000/docs",
        ],
        font_size=16,
        bold_first=True,
    )

    set_text(
        slide.shapes[12],
        [
            "Инструменты",
            "Google Colab: EDA, обучение, сравнение, интерпретация.",
            "FastAPI: REST API для инференса.",
            "Docker Compose: backend + frontend одной командой.",
            "Frontend: http://127.0.0.1:8080",
        ],
        font_size=16,
        bold_first=True,
    )

    add_card(
        slide,
        Inches(0.96),
        Inches(5.05),
        Inches(9.9),
        Inches(0.7),
        "Запуск демонстрации",
        ["docker compose up --build  →  frontend на :8080, backend на :8000"],
        accent=GREEN,
    )


def update_slide_11(prs: Presentation) -> None:
    slide = prs.slides[10]
    set_text(slide.shapes[2], ["Ограничения и развитие проекта"], font_size=28, bold_first=True)

    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            delete_shape(shape)

    add_card(
        slide,
        Inches(0.78),
        Inches(1.42),
        Inches(3.65),
        Inches(4.25),
        "Ограничения",
        [
            "Данные исторические: поведение пользователей меняется.",
            "Вероятность покупки не гарантирует покупку.",
            "Датасет может отражать сезонность и особенности источников трафика.",
        ],
        accent=CYAN,
    )
    add_card(
        slide,
        Inches(4.68),
        Inches(1.42),
        Inches(3.65),
        Inches(4.25),
        "Риски модели",
        [
            "Random Forest имеет признаки переобучения.",
            "False Negative — риск пропустить потенциального покупателя.",
            "False Positive — риск лишнего маркетингового действия.",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(8.58),
        Inches(1.42),
        Inches(3.65),
        Inches(4.25),
        "Как улучшить",
        [
            "Расширить поиск гиперпараметров.",
            "Подобрать бизнес-порог вероятности.",
            "Добавить мониторинг качества и регулярное переобучение.",
            "Собирать новые данные после изменений сайта.",
        ],
        accent=CYAN,
    )

    footer = slide.shapes.add_textbox(Inches(1.0), Inches(6.05), Inches(10.6), Inches(0.35))
    set_center_text(
        footer,
        "Вывод: решение уже пригодно для учебной демонстрации, а дальнейшее развитие связано с мониторингом и регулярным обновлением модели.",
        font_size=13,
        bold=False,
    )


def main() -> None:
    prs = Presentation(PPTX)
    update_slide_7(prs)
    update_slide_8(prs)
    update_slide_9(prs)
    update_slide_11(prs)
    prs.save(PPTX)

    if DOWNLOADS_COPY.parent.exists():
        copy2(PPTX, DOWNLOADS_COPY)

    print(PPTX)


if __name__ == "__main__":
    main()
