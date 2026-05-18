from __future__ import annotations

from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "Презентация_Покупательское_намерение_Матвей.pptx"
DOWNLOADS_COPY = (
    Path.home()
    / "Downloads"
    / "Презентация_Покупательское_намерение_Матвей.pptx"
)

DARK = RGBColor(67, 62, 76)
GREEN = RGBColor(137, 196, 0)
CYAN = RGBColor(50, 194, 215)
LIGHT = RGBColor(248, 251, 252)
MID = RGBColor(232, 240, 242)
TEXT = RGBColor(45, 45, 52)
WHITE = RGBColor(255, 255, 255)


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def set_text(
    shape,
    text: str,
    font_size: int,
    *,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.CENTER,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_cell(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    fill: RGBColor,
    line: RGBColor = MID,
    font_size: int = 15,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.CENTER,
) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    set_text(shape, text, font_size, bold=bold, color=color, align=align)


def add_note(slide, left, top, width, height, title: str, body: str, accent: RGBColor) -> None:
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.1)

    title_box = slide.shapes.add_textbox(left + Inches(0.14), top + Inches(0.08), width - Inches(0.28), Inches(0.28))
    set_text(title_box, title, 15, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    body_box = slide.shapes.add_textbox(left + Inches(0.14), top + Inches(0.42), width - Inches(0.28), height - Inches(0.48))
    set_text(body_box, body, 13, color=TEXT, align=PP_ALIGN.LEFT)


def redesign_slide_6(prs: Presentation) -> None:
    slide = prs.slides[5]

    # Keep background, title and slide number. Remove old chart/table/callout.
    for shape in list(slide.shapes):
        if shape.shape_id not in {
            slide.shapes[0].shape_id,
            slide.shapes[1].shape_id,
            slide.shapes[2].shape_id,
            slide.shapes[3].shape_id,
        }:
            delete_shape(shape)

    title = slide.shapes[2]
    set_text(title, "Сравнение моделей и выбор финальной", 27, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    # Context chips.
    add_cell(
        slide,
        Inches(0.88),
        Inches(1.28),
        Inches(3.0),
        Inches(0.42),
        "Основная метрика: F1-score",
        fill=GREEN,
        line=GREEN,
        font_size=14,
        bold=True,
        color=WHITE,
    )
    add_cell(
        slide,
        Inches(4.08),
        Inches(1.28),
        Inches(3.35),
        Inches(0.42),
        "Подбор: RandomizedSearchCV",
        fill=DARK,
        line=DARK,
        font_size=14,
        bold=True,
        color=WHITE,
    )
    add_cell(
        slide,
        Inches(7.62),
        Inches(1.28),
        Inches(3.75),
        Inches(0.42),
        "Дисбаланс: покупок около 15.47%",
        fill=CYAN,
        line=CYAN,
        font_size=14,
        bold=True,
        color=WHITE,
    )

    headers = ["Модель", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC"]
    rows = [
        ["Random Forest Tuned", "0.893", "0.648", "0.675", "0.662", "0.923"],
        ["CatBoost", "0.902", "0.727", "0.586", "0.649", "0.929"],
        ["Logistic Regression", "0.859", "0.529", "0.809", "0.640", "0.920"],
    ]

    x0 = Inches(0.86)
    y0 = Inches(1.95)
    widths = [Inches(3.05), Inches(1.45), Inches(1.55), Inches(1.45), Inches(1.2), Inches(1.55)]
    row_h = Inches(0.56)

    x = x0
    for header, width in zip(headers, widths, strict=True):
        add_cell(
            slide,
            x,
            y0,
            width,
            row_h,
            header,
            fill=DARK,
            line=WHITE,
            font_size=14,
            bold=True,
            color=WHITE,
        )
        x += width

    for row_index, row in enumerate(rows, start=1):
        y = y0 + row_h * row_index
        is_best = row_index == 1
        fill = RGBColor(235, 246, 222) if is_best else WHITE
        border = GREEN if is_best else MID
        x = x0
        for col_index, (value, width) in enumerate(zip(row, widths, strict=True)):
            add_cell(
                slide,
                x,
                y,
                width,
                row_h,
                value,
                fill=fill,
                line=border,
                font_size=15 if col_index == 0 else 16,
                bold=is_best or col_index == 0,
                color=DARK if is_best else TEXT,
                align=PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.CENTER,
            )
            x += width

    add_note(
        slide,
        Inches(0.88),
        Inches(4.35),
        Inches(3.55),
        Inches(1.25),
        "Почему выбрана RF Tuned",
        "Лучший F1-score: модель лучше балансирует поиск покупок и количество ложных срабатываний.",
        GREEN,
    )
    add_note(
        slide,
        Inches(4.68),
        Inches(4.35),
        Inches(3.55),
        Inches(1.25),
        "Почему не CatBoost",
        "ROC-AUC немного выше, но recall и F1 ниже основной модели, поэтому он не выбран финальным.",
        CYAN,
    )
    add_note(
        slide,
        Inches(8.48),
        Inches(4.35),
        Inches(3.05),
        Inches(1.25),
        "Baseline",
        "Logistic Regression даёт высокий recall, но больше ложных срабатываний.",
        DARK,
    )

    footer = slide.shapes.add_textbox(Inches(1.0), Inches(6.02), Inches(10.6), Inches(0.35))
    set_text(
        footer,
        "Вывод: финальная модель выбрана по F1-score, потому что задача несбалансирована и важно учитывать и precision, и recall.",
        13,
        color=DARK,
    )


def main() -> None:
    prs = Presentation(PPTX)
    redesign_slide_6(prs)
    prs.save(PPTX)

    if DOWNLOADS_COPY.parent.exists():
        copy2(PPTX, DOWNLOADS_COPY)

    print(PPTX)


if __name__ == "__main__":
    main()
