from __future__ import annotations

from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "Презентация_Покупательское_намерение_Матвей.pptx"
DOWNLOADS_COPY = (
    Path.home()
    / "Downloads"
    / "Презентация_Покупательское_намерение_Матвей.pptx"
)

DARK = RGBColor(67, 62, 76)
GREEN = RGBColor(137, 196, 0)
CYAN = RGBColor(50, 194, 215)
LIGHT = RGBColor(245, 248, 250)
TEXT = RGBColor(46, 46, 52)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, lines: list[str], font_size: int = 16, bold_first: bool = False) -> None:
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(4)
        if index == 0 and bold_first:
            paragraph.font.bold = True
            paragraph.font.color.rgb = DARK


def set_center_text(shape, text: str, font_size: int = 18, color=WHITE) -> None:
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color


def add_title(slide, title: str, number: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.52), Inches(9.7), Inches(0.55))
    set_text(title_box, [title], font_size=28, bold_first=True)

    line = slide.shapes.add_shape(1, Inches(0.22), Inches(1.34), Inches(12.55), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.color.rgb = GREEN

    num = slide.shapes.add_shape(1, Inches(12.0), Inches(6.35), Inches(0.42), Inches(0.33))
    num.fill.solid()
    num.fill.fore_color.rgb = GREEN
    num.line.color.rgb = GREEN
    set_center_text(num, number, font_size=11)


def add_card(slide, left, top, width, height, title: str, body: list[str], accent=GREEN) -> None:
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), Inches(0.35))
    set_text(title_box, [title], font_size=17, bold_first=True)

    body_box = slide.shapes.add_textbox(left + Inches(0.16), top + Inches(0.58), width - Inches(0.32), height - Inches(0.7))
    set_text(body_box, body, font_size=14)


def add_pipeline(slide) -> None:
    labels = ["Обучение", "MLflow", "DagsHub", "Версии модели"]
    lefts = [0.9, 3.75, 6.6, 9.45]
    for left, label in zip(lefts, labels, strict=True):
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.52), Inches(2.25), Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = GREEN if label in {"MLflow", "DagsHub"} else DARK
        box.line.color.rgb = box.fill.fore_color.rgb
        set_center_text(box, label, font_size=16)

    for left in [3.25, 6.1, 8.95]:
        arrow = slide.shapes.add_textbox(Inches(left), Inches(1.63), Inches(0.35), Inches(0.35))
        set_text(arrow, ["→"], font_size=23, bold_first=True)


def move_last_slide_before_previous(prs: Presentation) -> None:
    """Move a newly appended slide before the previous last slide."""
    sld_id_lst = prs.slides._sldIdLst
    slide_ids = list(sld_id_lst)
    if len(slide_ids) < 2:
        return
    new_slide_id = slide_ids[-1]
    previous_last_id = slide_ids[-2]
    sld_id_lst.remove(new_slide_id)
    previous_index = list(sld_id_lst).index(previous_last_id)
    sld_id_lst.insert(previous_index, new_slide_id)


def main() -> None:
    prs = Presentation(PPTX)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "MLflow + DagsHub: управление экспериментами", "12")
    add_pipeline(slide)

    add_card(
        slide,
        Inches(0.8),
        Inches(2.6),
        Inches(3.5),
        Inches(2.75),
        "Что логируется",
        [
            "Параметры модели и random_state.",
            "Метрики: Accuracy, Precision, Recall, F1, ROC-AUC.",
            "Артефакты: model_comparison.csv, metadata.json, confusion matrix.",
            "Финальный sklearn pipeline.",
        ],
        accent=CYAN,
    )
    add_card(
        slide,
        Inches(4.7),
        Inches(2.6),
        Inches(3.5),
        Inches(2.75),
        "Зачем это нужно",
        [
            "Результаты не теряются после перезапуска Colab.",
            "Можно сравнивать разные версии моделей.",
            "Видно, почему выбрана финальная модель.",
            "Проще воспроизвести эксперимент.",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(8.6),
        Inches(2.6),
        Inches(3.5),
        Inches(2.75),
        "MLOps-эффект",
        [
            "DagsHub выступает удалённым MLflow Tracking Server.",
            "Фиксируется жизненный цикл модели: запуск, метрики, артефакты.",
            "Это база для мониторинга и регулярного переобучения.",
        ],
        accent=CYAN,
    )

    footer = slide.shapes.add_textbox(Inches(1.0), Inches(5.95), Inches(10.8), Inches(0.42))
    set_text(
        footer,
        ["Вывод: MLflow + DagsHub добавляют проекту воспроизводимость и показывают базовый подход к управлению ML-экспериментами."],
        font_size=13,
    )

    move_last_slide_before_previous(prs)
    prs.save(PPTX)

    if DOWNLOADS_COPY.parent.exists():
        copy2(PPTX, DOWNLOADS_COPY)

    print(PPTX)


if __name__ == "__main__":
    main()
